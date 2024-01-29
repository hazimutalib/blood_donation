import requests
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from spire.presentation import Presentation as Presentation2, FileFormat
import base64
from pptx.util import Inches

def set_font_properties(text_frame, font_name, font_size, font_color, is_bold, alignment):
    if text_frame is not None:

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*font_color)
                run.font.bold = is_bold

            paragraph.alignment = alignment



def edit_powerpoint_template(template_path, output_path, malaysia, kuala_lumpur, kedah, perak, johor, sarawak, pulau_pinang, sabah, melaka, selangor, 
                             negeri_sembilan, terengganu, pahang, kelantan, max_date, daily_malaysia, daily_kuala_lumpur, daily_kedah, daily_perak, daily_johor,
                             daily_sarawak, daily_pulau_pinang, daily_sabah, daily_melaka, daily_selangor, 
                             daily_negeri_sembilan, daily_terengganu, daily_pahang, daily_kelantan):
    # Load the template presentation
    presentation = Presentation(template_path)

    # Font properties for the title and content
    content_font_name = "Aharoni"
    content_font_size = 56
    content_font_color = (0, 0, 0)  # RGB color tuple, e.g., black
    content_is_bold = True
    content_alignment = PP_ALIGN.CENTER

    # Access and modify slides
    for i, slide in enumerate(presentation.slides):

        # If there is a placeholder named "Content", update it
        for shape in slide.shapes:
            # picture = slide.shapes.add_picture('C:/Users/Analyst07/Desktop/blood_donation/test.png', Inches(2), Inches(1), Inches(3), Inches(3))
            if shape.has_text_frame:
                
                if "malaysia" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("malaysia", "{:,}".format(malaysia))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, (255, 255, 255), content_is_bold, content_alignment)

                if "perlis" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("perlis", "0")
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "kedah" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("kedah", "{:,}".format(kedah))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)
                
                if "pulau pinang" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("pulau pinang", "{:,}".format(pulau_pinang))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "perak" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("perak", "{:,}".format(perak))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "selangor" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("selangor", "{:,}".format(selangor))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "kuala lumpur" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("kuala lumpur", "{:,}".format(kuala_lumpur))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "putrajaya" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("putrajaya", "0")
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "n. sembilan" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("n. sembilan", "{:,}".format(negeri_sembilan))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)
                
                if "sabah" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("sabah", "{:,}".format(sabah))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "labuan" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("labuan", "0")
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)
                
                if "sarawak" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("sarawak", "{:,}".format(sarawak))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "kelantan" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("kelantan", "{:,}".format(kelantan))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "terengganu" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("terengganu", "{:,}".format(terengganu))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "pahang" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("pahang", "{:,}".format(pahang))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "melaka" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("melaka", "{:,}".format(melaka))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "johor" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("johor", "{:,}".format(johor))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "max_date" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("max_date", "{}".format(max_date))
                    set_font_properties(shape.text_frame, 'Ahoroni', 18, content_font_color, content_is_bold, content_alignment)

                if "last_date" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("last_date", "{}".format(max_date))
                    set_font_properties(shape.text_frame, 'Verdana', 48, content_font_color, content_is_bold, content_alignment)

                if "daily_mly" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_mly", "{:,}".format(daily_malaysia))
                    set_font_properties(shape.text_frame, content_font_name, 24, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_pls" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_pls", "0")
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_kdh" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_kdh", "{:,}".format(daily_kedah))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)
                
                if "daily_png" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_png", "{:,}".format(daily_pulau_pinang))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_prk" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_prk", "{:,}".format(daily_perak))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_slgr" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_slgr", "{:,}".format(daily_selangor))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_kl" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_kl", "{:,}".format(daily_kuala_lumpur))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_pjy" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_pjy", "0")
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_ns" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_ns", "{:,}".format(daily_negeri_sembilan))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)
                
                if "daily_sbh" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_sbh", "{:,}".format(daily_sabah))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_lbn" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_lbn", "0")
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)
                
                if "daily_srwk" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_srwk", "{:,}".format(daily_sarawak))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_ktn" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_ktn", "{:,}".format(daily_kelantan))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_trgn" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_trgn", "{:,}".format(daily_terengganu))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_phg" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_phg", "{:,}".format(daily_pahang))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_mlk" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_mlk", "{:,}".format(daily_melaka))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)

                if "daily_jhr" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("daily_jhr", "{:,}".format(daily_johor))
                    set_font_properties(shape.text_frame, content_font_name, 20, (28, 182, 54), content_is_bold, content_alignment)


    # Save the modified presentation
    presentation.save(output_path)

    print(f"Edited presentation saved to {output_path}")


