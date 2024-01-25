import requests
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from spire.presentation import Presentation as Presentation2, FileFormat
import base64


def set_font_properties(text_frame, font_name, font_size, font_color, is_bold, alignment):
    if text_frame is not None:

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*font_color)
                run.font.bold = is_bold

            paragraph.alignment = alignment



def edit_powerpoint_template(template_path, output_path, malaysia, kuala_lumpur, kedah, perak, johor, sarawak, pulau_pinang, sabah, melaka, selangor, 
                             negeri_sembilan, terengganu, pahang, kelantan, max_date):
    # Load the template presentation
    presentation = Presentation(template_path)

    # Font properties for the title and content
    content_font_name = "Verdana"
    content_font_size = 56
    content_font_color = (0, 0, 0)  # RGB color tuple, e.g., black
    content_is_bold = True
    content_alignment = PP_ALIGN.CENTER

    # Access and modify slides
    for i, slide in enumerate(presentation.slides):

        # If there is a placeholder named "Content", update it
        for shape in slide.shapes:
            if shape.has_text_frame:
                
                if "malaysia" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("malaysia", "{:,}".format(malaysia))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, (255, 255, 255), content_is_bold, content_alignment)

                if "perlis" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("perlis", "0")
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "kedah" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("kedah", "{:,}".format(kedah))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)
                
                if "pulau pinang" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("pulau pinang", "{:,}".format(pulau_pinang))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "perak" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("perak", "{:,}".format(perak))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "selangor" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("selangor", "{:,}".format(selangor))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "kuala lumpur" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("kuala lumpur", "{:,}".format(kuala_lumpur))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "putrajaya" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("putrajaya", "0")
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "n. sembilan" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("n. sembilan", "{:,}".format(negeri_sembilan))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)
                
                if "sabah" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("sabah", "{:,}".format(sabah))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "labuan" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("labuan", "0")
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)
                
                if "sarawak" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("sarawak", "{:,}".format(sarawak))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "kelantan" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("kelantan", "{:,}".format(kelantan))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "terengganu" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("terengganu", "{:,}".format(terengganu))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "pahang" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("pahang", "{:,}".format(pahang))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "melaka" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("melaka", "{:,}".format(melaka))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "johor" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("johor", "{:,}".format(johor))
                    set_font_properties(shape.text_frame, content_font_name, content_font_size, content_font_color, content_is_bold, content_alignment)

                if "date" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("date", "{}".format(max_date))
                    set_font_properties(shape.text_frame, 'Verdana', 16, content_font_color, content_is_bold, content_alignment)

                if "count" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("count", "{}".format(max_date))
                    set_font_properties(shape.text_frame, 'Verdana', 46, content_font_color, content_is_bold, content_alignment)


    # Save the modified presentation
    presentation.save(output_path)

    print(f"Edited presentation saved to {output_path}")



def upload_pptx_to_github(repo_owner, repo_name, template_path, file_path, github_token, malaysia, kuala_lumpur, kedah, perak, johor, sarawak, pulau_pinang, sabah, melaka, selangor, 
                             negeri_sembilan, terengganu, pahang, kelantan):
    

    edit_powerpoint_template(template_path, file_path, malaysia, kuala_lumpur, kedah, perak, johor, sarawak, pulau_pinang, sabah, melaka, selangor, 
                             negeri_sembilan, terengganu, pahang, kelantan)


    with open(file_path, 'rb') as file:
        file_content = base64.b64encode(file.read()).decode('utf-8')


    api_url = f'https://api.github.com/repos/{repo_owner}/{repo_name}/contents/{file_path}'

    # Check if the file exists
    response = requests.get(api_url, headers={'Authorization': f'Token {github_token}'})
    if response.status_code == 200:
        # File exists, update it
        sha = response.json()['sha']
    else:
        # File doesn't exist, create it
        sha = None

    # Create a commit with the updated presentation file
    commit_message = 'Automated commit: Add/update PowerPoint presentation'
    commit_data = {
        'message': commit_message,
        'content': file_content,
        'sha': sha  # Include the SHA here
    }

    response = requests.put(api_url, headers={'Authorization': f'Token {github_token}'}, json=commit_data)

    if response.status_code == 200:
        print('Presentation file successfully saved to GitHub.')
    else:
        print(f'Failed to save presentation file. Status code: {response.status_code}, Message: {response.text}')



def upload_pdf_to_github(pdf_file_path, github_token, repo_owner, repo_name, commit_message="Add PDF file"):
    # Set up the API endpoint
    url = f'https://api.github.com/repos/{repo_owner}/{repo_name}/contents/{pdf_file_path}'

    # Read the PDF file content
    with open(pdf_file_path, 'rb') as pdf_file:
        pdf_content = pdf_file.read()

    # Encode the content in base64
    encoded_content = base64.b64encode(pdf_content).decode('utf-8')

    # Create the API payload
    payload = {
        'message': commit_message,
        'content': encoded_content,
    }

    # Set the Authorization header with the GitHub token
    headers = {
        'Authorization': f'token {github_token}',
    }

    # Make the API request to create or update the file
    response = requests.put(url, json=payload, headers=headers)

    if response.status_code == 201:
        print(f'PDF file {pdf_file_path} successfully uploaded to GitHub.')
    else:
        print(f'Failed to upload PDF file. Status code: {response.status_code}\nResponse content: {response.text}')






# pdf_file_path = "path/to/your/pdf_file.pdf"
# github_token = "your_github_token"
# repo_owner = "your_username_or_organization"
# repo_name = "your_repository_name"

# upload_pdf_to_github(pdf_file_path, github_token, repo_owner, repo_name)



# repo_owner = 'hazimutalib'
# repo_name = 'blood_donation'
# template_path = './blood_donation.pptx'
# file_path = './infographic/output_test.pptx'
# github_token = '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'
# upload_pptx_to_github(repo_owner, repo_name, template_path, file_path, github_token)

# file_path_pdf = './infographic/output_test.pdf'
# presentation2 = Presentation2()
# presentation2.LoadFromFile(file_path)
# presentation2.SaveToFile(file_path_pdf, FileFormat.PDF)
# # presentation2.Dispose()

# upload_pptx_to_github(repo_owner, repo_name, template_path, file_path_pdf, github_token)